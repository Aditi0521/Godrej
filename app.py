import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE  # Import MSO_SHAPE_TYPE
from docx import Document
import streamlit as st
from langchain.text_splitter import CharacterTextSplitter
from cohere import Client as CohereClient
from fpdf import FPDF
import tempfile
import fitz
from PIL import Image, ImageDraw
import pytesseract
import io
import base64
from gtts import gTTS  
from deep_translator import GoogleTranslator  

def load_cohere_api_key():
    try:
        api_key = os.environ["COHERE_API_KEY"]
        return api_key
    except KeyError as e:
        raise KeyError(f"KeyError: {e}. Ensure 'COHERE_API_KEY' is added to secrets.")

def process_text(text, chunk_size, chunk_overlap):
    if chunk_size is None:
        chunk_size = len(text)
    if chunk_overlap is None:
        chunk_overlap = 0

    if not isinstance(chunk_overlap, int) or not isinstance(chunk_size, int):
        raise TypeError("chunk_overlap and chunk_size must be integers")

    if chunk_overlap > chunk_size:
        raise ValueError("chunk_overlap cannot be greater than chunk_size")

    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def extract_text_from_pdf(pdf):
    pdf.seek(0)
    pdf_document = fitz.open(stream=pdf.read(), filetype="pdf")
    text = ""
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    return text

def extract_text_from_pptx(pptx_file):
    text = []
    presentation = Presentation(pptx_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_images_from_pptx(file):
    prs = Presentation(file)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image.blob
                img_io = io.BytesIO(img)
                images.append(Image.open(img_io))
    return images

def extract_images_from_pdf(pdf):
    pdf.seek(0)
    images = []
    pdf_document = fitz.open(stream=pdf.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            images.append(image_bytes)
    return images

def extract_text_from_docx(docx_file):
    text = []
    doc = Document(docx_file)
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)

def extract_images_from_docx(file):
    doc = Document(file)
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image = Image.open(io.BytesIO(rel.target_part.blob))
            images.append(image)
    return images

def get_image_descriptions(images):
    descriptions = []
    for img in images:
        try:
            text = pytesseract.image_to_string(img)
            descriptions.append(text)
        except Exception as e:
            descriptions.append(f"Error processing image: {e}")
    return descriptions

def summarize_text(text, image_descriptions, co):
    if image_descriptions:
        full_text = text + "\n\n" + "\n".join(image_descriptions)
    else:
        full_text = text
    text_length = len(full_text)

    if text_length > 50000:
        chunk_size = 5000
        chunk_overlap = 1000
        max_tokens = 1000
        summary_prompt = "Summarize the following text concisely:"
    elif text_length > 10000:
        chunk_size = 2000
        chunk_overlap = 500
        max_tokens = 500
        summary_prompt = "Summarize the following key points:"
    else:
        chunk_size = 1000
        chunk_overlap = 200
        max_tokens = 300
        summary_prompt = "Provide a concise summary of the following text:"

    chunks = process_text(full_text, chunk_size, chunk_overlap)
    summary_text = ""
    for chunk in chunks:
        response = co.generate(
            prompt=summary_prompt + chunk,
            model="command-xlarge-nightly",
            max_tokens=max_tokens,
            temperature=0.7,
        )
        summary_text += response.generations[0].text.strip() + "\n\n"

    # Enhance readability
    summary_lines = summary_text.split('\n')
    structured_summary = ""
    for line in summary_lines:
        if len(line) > 100:
            structured_summary += "\n" + line.strip() + "\n"
        elif len(line) > 0:
            structured_summary += "\n- " + line.strip()

    return structured_summary

def translate_text(text, target_language):
    translator = GoogleTranslator(source='auto', target=target_language)
    translation = translator.translate(text)
    return translation

def answer_question(question, text, co):
    response = co.generate(
        prompt=question + "\n\n" + text,
        model="command-xlarge-nightly",
        max_tokens=200,
        temperature=0.7,
    )
    return response.generations[0].text.strip()

class PDF(FPDF):
    def header(self):
        if hasattr(self, 'header_text'):
            self.set_font('DejaVu', '', 14)
            self.cell(0, 10, self.header_text, 0, 1, 'L')

    def chapter_title(self, title):
        self.set_font('DejaVu', '', 12)
        self.cell(0, 10, title, 0, 1, 'L')
        self.ln(10)

    def chapter_body(self, body):
        self.set_font('DejaVu', '', 12)
        self.multi_cell(0, 10, body)
        self.ln()

def generate_pdf_report(summary_text, translated_summary, qna_history):
    pdf = PDF()

    font_path = os.path.join(os.path.dirname(__file__), 'DejaVuSans.ttf')
    pdf.add_font('DejaVu', '', font_path, uni=True)
    pdf.set_font('DejaVu', '', 12)

    pdf.add_page()
    pdf.set_left_margin(10)
    pdf.set_right_margin(10)
    pdf.set_xy(10, 30)
    pdf.chapter_title("Summary:")
    pdf.chapter_body(summary_text if summary_text else "No summary available.")

    if translated_summary and translated_summary != summary_text:
        pdf.add_page()
        pdf.header_text = "Translated Summary"
        pdf.header()
        pdf.chapter_body(translated_summary)

    if qna_history:
        pdf.add_page()
        pdf.header_text = "Q&A History"
        pdf.header()
        for qna in qna_history:
            pdf.chapter_body(f"Q: {qna[0]}\nA: {qna[1]}")

    pdf_file_path = 'translated_summary.pdf'
    pdf.output(pdf_file_path, 'F')

    return pdf_file_path

def display_pdf(file):
    file.seek(0)
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        st.image(img, caption=f'Page {page_num+1}', use_column_width=True)

def display_docx_as_pdf(file):
    # Extract text from the DOCX file
    text = extract_text_from_docx(file)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as pdf_file:
        pdf = PDF()
        font_path = os.path.join(os.path.dirname(__file__), 'DejaVuSans.ttf')
        pdf.add_font('DejaVu', '', font_path, uni=True)
        pdf.set_font("DejaVu", size=12)
        pdf.add_page()
        pdf.multi_cell(0, 10, text)
        pdf.output(pdf_file.name)
        
        pdf_file.seek(0)
        pdf_data = pdf_file.read()
    
    st.download_button("Download PDF", pdf_data, file_name="document.pdf", mime="application/pdf")

def main():
    st.set_page_config(page_title="Document Summarizer and Q&A", layout="wide")

    # Initialize session state
    if 'qna_history' not in st.session_state:
        st.session_state.qna_history = []

    col1, col2 = st.columns(2)
    with col1:
        st.header("Uploaded Document")
        doc_type = st.selectbox("Select Document Type:", ("PDF", "PPT", "Word"))
        uploaded_file = st.file_uploader("Upload Your Document:")

        if uploaded_file is not None:
            if doc_type == "PDF":
                display_pdf(uploaded_file)
            elif doc_type == "Word":
                display_docx_as_pdf(uploaded_file)
            uploaded_file.seek(0)

    with col2:
        st.header("Summary and Q&A")
        if uploaded_file is not None:
            cohere_api_key = load_cohere_api_key()
            co = CohereClient(cohere_api_key)
            if doc_type == "PDF":
                text = extract_text_from_pdf(uploaded_file)
                images = extract_images_from_pdf(uploaded_file)
            elif doc_type == "PPT":
                text = extract_text_from_pptx(uploaded_file)
                images = extract_images_from_pptx(uploaded_file)
            elif doc_type == "Word":
                text = extract_text_from_docx(uploaded_file)
                images = extract_images_from_docx(uploaded_file)
            image_descriptions = get_image_descriptions(images)
            summary_text = summarize_text(text, image_descriptions, co)
            st.subheader("Document Summary")
            st.write(summary_text)
            target_language = st.selectbox("Select Language for Translation:", ["en", "es", "fr", "de", "it", "hi"])
            if target_language != "en":
                translated_summary = translate_text(summary_text, target_language)
                st.subheader(f"Translated Summary ({target_language}):")
                st.write(translated_summary)
            else:
                translated_summary = summary_text

            st.subheader("Ask a Question")
            question = st.text_input("Enter your question:")
            if question:
                answer = answer_question(question, text, co)
                st.write(f"Answer: {answer}")
                st.session_state.qna_history.append((question, answer))

            if st.session_state.qna_history:
                st.subheader("Q&A History")
                for q, a in st.session_state.qna_history[:-1]:  # Display all previous Q&A except the current one
                    st.write(f"Q: {q}")
                    st.write(f"A: {a}")

            if question:
                st.subheader("Current Q&A")
                st.write(f"Q: {question}")
                st.write(f"A: {answer}")

            if st.button("Download Report"):
                pdf_file_path = generate_pdf_report(summary_text, translated_summary, st.session_state.qna_history)
                with open(pdf_file_path, "rb") as pdf_file:
                    b64_pdf = base64.b64encode(pdf_file.read()).decode("utf-8")
                    pdf_link = f'<a href="data:application/octet-stream;base64,{b64_pdf}" download="report.pdf">Download PDF Report</a>'
                    st.markdown(pdf_link, unsafe_allow_html=True)

            if st.button("Text to Speech"):
                tts = gTTS(text=summary_text, lang='en')
                with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_audio_file:
                    tts.save(temp_audio_file.name)
                    temp_audio_file_path = temp_audio_file.name
                audio_file = open(temp_audio_file_path, "rb")
                audio_bytes = audio_file.read()
                st.audio(audio_bytes, format='audio/mp3')
                os.remove(temp_audio_file_path)

if __name__ == "__main__":
    main()
